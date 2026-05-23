"""Generate IMT542_Final_Project.pptx for class presentation."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "presentation" / "IMT542_Final_Project.pptx"
UW_PURPLE = RGBColor(0x4B, 0x2E, 0x83)

SLIDES = [
    (
        "Neighborhood Small Business Platform",
        "IMT 542 Final Project — Portable Information Structures\n"
        "Sneha Reddy | Spring 2026\n\n"
        "GitHub: github.com/Sneha-8/IMT-542/tree/main/final",
    ),
    (
        "1. Information Story",
        "User: Seattle resident who wants to discover and support local shops\n\n"
        "Problem: Business information is scattered across Yelp, HTML directories, "
        "and social media — not portable for apps or civic tools\n\n"
        "Outcome: One machine-readable JSON feed with provenance and quality metrics\n\n"
        "Insight: Visualize listings + analyze relationships (stats, open-now)",
    ),
    (
        "Users and Requirements",
        "Residents — search by neighborhood, tags, open now\n"
        "Owners — stable public listing; restricted contact fields\n"
        "Developers — JSON Schema contract + provenance\n\n"
        "In scope: SBP v1.0 API, FAIR metadata, NoSQL, auth-gated contact\n"
        "Out of scope: payments, Yelp production key in repo, native mobile app",
    ),
    (
        "2. Existing Structures and FAIR",
        "Sources analyzed:\n"
        "  - Yelp Fusion API (nested vendor JSON, license limits)\n"
        "  - HTML neighborhood directories (no stable IDs)\n"
        "  - I7 static JSON (no provenance or schema version)\n\n"
        "FAIR gaps: weak reusability, inconsistent access, no completeness metrics\n"
        "SBP adds: stable IDs, HTTPS REST, JSON Schema, provenance.license",
    ),
    (
        "3. Structure Improvements",
        "At least 4 dimensions changed:\n"
        "  Information — quality_flags, completeness_score\n"
        "  Structure — SBP envelope + flattened LocalBusiness fields\n"
        "  Format — schema_version 1.0 + sbp_schema_v1.json\n"
        "  Access — REST /businesses/search vs file-in-memory (I7)",
    ),
    (
        "New Portable Structure",
        "Top-level: schema_version, provenance, query_summary, businesses[]\n\n"
        "Per business: id, name, categories, rating, location, contact, "
        "tags, hours, products, quality_flags\n\n"
        "Normalizer: small_businesses.json -> SBP record (normalizer.py)",
    ),
    (
        "4. Functional System",
        "Stack: JSON source -> seed_db.py -> MongoDB/Mongita -> Flask API\n\n"
        "Run locally:\n"
        "  export USE_MONGITA=true && python seed_db.py\n"
        "  flask --app app run -p 5002\n"
        "  ngrok http 5002  (submit HTTPS URL to Canvas)\n\n"
        "Demo client: python access_api.py",
    ),
    (
        "Live Demo Endpoints",
        "GET /health\n"
        "GET /schema\n"
        "GET /businesses/search?term=coffee&location=Capitol Hill\n"
        "GET /businesses/sb-001\n"
        "GET /stats  |  GET /open-now\n\n"
        "HTML slides: final/presentation/index.html?api=<ngrok-url>",
    ),
    (
        "Quality and Performance",
        "Desired: completeness >= 0.85, schema validates 100%\n"
        "Observed: ~0.92 completeness on demo set; 9 pytest tests pass\n"
        "Latency: warm search < 50 ms local\n\n"
        "Remediation: Yelp live adapter, UptimeRobot, Locust load tests (G9)",
    ),
    (
        "Security and Ethics",
        "Bearer token unlocks email/phone (SBP_API_TOKEN)\n"
        "Public responses strip restricted contact fields\n"
        "No secrets in repo; ethics doc in final/docs/06_Ethics_and_Limitations.md\n\n"
        "Aligns with G9 test plan",
    ),
    (
        "Thank You / Questions",
        "Repository: github.com/Sneha-8/IMT-542\n"
        "Documentation: final/docs/\n"
        "Presentation HTML: final/presentation/index.html\n\n"
        "Questions?",
    ),
]


def add_slide(prs, title: str, body: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = UW_PURPLE

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(9), Inches(5.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(18)
        para.space_after = Pt(6)


def main():
    prs = Presentation()
    for title, body in SLIDES:
        add_slide(prs, title, body)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
